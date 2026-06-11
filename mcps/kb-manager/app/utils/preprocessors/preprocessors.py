import re 
import hashlib
import json
from pathlib import Path
from typing import Optional
import pandas as pd
from typing import List, Optional, Set
import pdfplumber
from docx import Document
from pptx import Presentation
from dataclasses import dataclass, asdict
from utils.logger import setup_logger
# ------------------------------
# Регулярные паттерны для поиска
#  -----------------------------
HEADER_RE = re.compile(r"^(#{1,6})\s+(.*)") # паттерн поиска 
URL_PATTERN = re.compile(r"https?://[^\s\)\]]+") # паттерн ссылок
MD_RE = re.compile(r"[*_`#>~\-]+")
PUNCT_RE = re.compile(r"[^\w\s]")
MULTISPACE_RE = re.compile(r"\s+")
QUESTION_WORD_RE = re.compile(
    r"(?i)\b(" + "|".join([
        "где", "когда", "почему", "зачем", "кто", "какой", "какая", "какие",
        "можно ли", "обязательно ли", "нужно ли", "стоит ли", "возможно ли", "допустимо ли"
    ]) + r")\b"
)
QUESTION_WORDS = [
    "где", "когда", "почему", "зачем", "кто",
    "какой", "какая", "какие",
    "можно ли", "обязательно ли", "нужно ли", "стоит ли",
    "возможно ли", "допустимо ли"
]
# паттерн на вопросные слова
QUESTION_WORD_RE = re.compile(
    r"(?i)\b(" + "|".join(re.escape(w) for w in QUESTION_WORDS) + r")\b"
)
# патерны вопросов
QA_PATTERNS = [
    (re.compile(r"^\*\*Вопрос:\*\*\s*(.+)"), "Q"),
    (re.compile(r"^\*\*Ответ:\*\*\s*(.+)"), "A"),
    (re.compile(r"^Q:\s*(.+)"), "Q"),
    (re.compile(r"^A:\s*(.+)"), "A"),
]

MD_STARS = re.compile(r"\*\*")
QUESTION_SELECT = re.compile(r"вопрос[:\s]*")
ANSWER_SELECT = re.compile(r"ответ[:\s]*")

# -----------
# Data Models
# -----------
@dataclass
class RawFAQItem:
    """Промежуточная структура перед нормализацией"""
    question: str
    answer: str 
    section_path: List[str]
    source_file: str
    source_type: str # 'markdown', 'excel', 'csv' 

@dataclass
class NormalizedFAQItem:
    """Итоговая структура для JSON"""
    question: str
    answer: str
    category: str
    section_path: List[str]
    source_file: str
    source_type: str
    hash: str
    canonical_question: str
    canonical_answer: str

class FAQPreprocessor:
    def __init__(self, output_file: Path, log_dir: Path):
        self.logger = setup_logger("Preprocessor Parser", log_dir)
        self.output_file = Path(output_file)
        self.items: List[NormalizedFAQItem] = []
        self.seen_hashes: Set[str] = set()
        self.errors: List[str] = []

    # --- Методы помощники ---
    @staticmethod
    def canonical_text(text: str) -> str:
        """Очистка текста для хеширования и поиска"""
        if not isinstance(text, str):
            text = str(text)
        text = text.lower()
        text = MD_STARS.sub("", text)               # убрать markdown **
        text = QUESTION_SELECT.sub("", text)        # убрать "вопрос"
        text = ANSWER_SELECT.sub("", text)
        text = URL_PATTERN.sub("", text)
        text = MD_RE.sub("", text)  # Убираем markdown символы
        text = PUNCT_RE.sub(" ", text)
        text = MULTISPACE_RE.sub(" ", text)
        return text.strip()

    def _normalize_item(self, raw: RawFAQItem) -> Optional[NormalizedFAQItem]:
        """Превращает сырые данные в финальный формат с хешем"""
        q_clean = raw.question.strip()
        a_clean = raw.answer.strip()
        
        # Валидация
        if not q_clean or not a_clean:
            return None
        if len(q_clean) < 3: # Минимальная длина вопроса
            return None

        canonical_q = self.canonical_text(q_clean)
        canonical_a = self.canonical_text(a_clean)
        context = "|".join([
            canonical_q, 
            canonical_a,
            # raw.source_file,
            # "/".join(raw.section_path)
        ])
        # Генерация хеша
        item_hash = hashlib.sha256((context).encode("utf-8")).hexdigest()

        # Дедупликация
        if item_hash in self.seen_hashes:
            return None
        self.seen_hashes.add(item_hash)

        category = raw.section_path[0] if raw.section_path else "General"

        return NormalizedFAQItem(
            question=q_clean,
            answer=a_clean,
            category=category,
            section_path=raw.section_path,
            source_file=str(raw.source_file),
            source_type=raw.source_type,
            hash=item_hash,
            canonical_question=canonical_q,
            canonical_answer=canonical_a
        )

    # --- Parsers ---
    def _parse_markdown(self, file_path: Path) -> List[RawFAQItem]:
        results = []
        try:
            with open(file_path, encoding="utf-8") as f:
                lines = f.readlines()
        except Exception as e:
            self.logger.error(f"Error reading MD {file_path}: {e}")
            return []

        # Логика стека для заголовков (более надежная, чем фиксированные индексы)
        # stack хранит кортежи (level, title)
        header_stack: List[tuple[int, str]] = [] 
        
        current_question = None
        current_answer_buffer = []

        def get_current_path():
            return [title for _, title in header_stack]

        def flush_qa():
            nonlocal current_question, current_answer_buffer
            if current_question and current_answer_buffer:
                results.append(RawFAQItem(
                    question=current_question,
                    answer="\n".join(current_answer_buffer).strip(),
                    section_path=get_current_path(),
                    source_file=file_path.name,
                    source_type="markdown"
                ))
            current_question = None
            current_answer_buffer = []

        for line in lines:
            line = line.rstrip()
            header_match = HEADER_RE.match(line)

            # Если встретили заголовок
            if header_match:
                flush_qa() # Сохраняем предыдущую пару, если была
                level = len(header_match.group(1))
                title = header_match.group(2).strip()
                
                # Очищаем стек от заголовков того же уровня или глубже
                while header_stack and header_stack[-1][0] >= level:
                    header_stack.pop()
                header_stack.append((level, title))
                continue

            # Эвристика вопроса
            is_question = False
            clean_line = line.strip().lower()
            
            # Если строка начинается с маркеров вопроса
            if (clean_line.startswith("**вопрос") or 
                clean_line.startswith("q:") or 
                clean_line.startswith("вопрос:")):
                is_question = True
            # Или содержит '?' в конце (и не является ссылкой)
            elif "?" in clean_line and not URL_PATTERN.search(clean_line):
                # Доп проверка: не слишком ли длинный "вопрос" (чтобы не захватить абзац текста)
                if len(clean_line) < 300: 
                    is_question = True

            if is_question:
                flush_qa()
                # Чистим мусор (**, Q:, Вопрос:)
                cleaned_q = re.sub(r"^(\*\*|Q:|Вопрос:)\s*", "", line.strip(), flags=re.IGNORECASE)
                cleaned_q = re.sub(r"\*\*$", "", cleaned_q)
                current_question = cleaned_q.strip()
            elif current_question:
                # Если мы внутри ответа
                if line.strip(): # Пропускаем пустые строки, если нужно, или сохраняем форматирование
                    current_answer_buffer.append(line)

        flush_qa() # Сохраняем последний элемент
        return results

    def _load_tables(self, file_path: Path) -> dict[str, pd.DataFrame]:
        """Возвращает словарь сопоставления страниц в файле с dataframe"""
        if file_path.suffix == ".csv":
            # так как csv не может содержать такую информацию из-за своего формата
            df = pd.read_csv(file_path)
            return {"__default__": df}
        if file_path.suffix in [".xlsx", ".xls"]:
            return pd.read_excel(file_path, sheet_name=None)
        raise ValueError(f"Unsupported table format: {file_path}")
    
    def _parse_dataframe(self, df: pd.DataFrame, 
                         file_path: Path, 
                         sheet_name: Optional[str]=None
                         ) -> List[RawFAQItem]:
        results = []
        # Нормализация имен колонок (приводим к нижнему регистру для поиска)
        df.columns = [str(c).lower().strip() for c in df.columns]
        # Маппинг возможных названий колонок
        col_map = {
            'q': ['вопрос', 'question', 'q'],
            'a': ['ответ', 'answer', 'a'],
            'cat': ['раздел', 'категория', 'category', 'section']
        }
         # Поиск реальных колонок
        real_cols = {}
        for key, candidates in col_map.items():
            for cand in candidates:
                if cand in df.columns:
                    real_cols[key] = cand
                    break
        
        if 'q' not in real_cols or 'a' not in real_cols:
            self.logger.warning(f"Skipping {file_path}: Required columns (Question, Answer) not found.")
            return []

        df = df.dropna(subset=[real_cols['q'], real_cols['a']])
        
        for _, row in df.iterrows():
            category = (
                str(row[real_cols['cat']]).strip() 
                if 'cat' in real_cols and not pd.isna(row[real_cols['cat']])
                else "General"
            )
            # Для таблиц section_path делаем списком из одного элемента - категории
            section_path = []
            if sheet_name:
                section_path.append(sheet_name)
            if category and category not in section_path:
                section_path.append(category)
            
            results.append(RawFAQItem(
                question=str(row[real_cols['q']]),
                answer=str(row[real_cols['a']]),
                section_path=section_path,
                source_file=file_path.name,
                source_type=file_path.suffix.lstrip('.')
            ))
        return results

    def _parse_table(self, file_path: Path) -> List[RawFAQItem]:
        results = []
        try:
            tables = self._load_tables(file_path)
            for sheet_name, df in tables.items():
                self.logger.info(f"Processing table {file_path.name} / sheet={sheet_name}")
                results.extend(
                    self._parse_dataframe(
                        df=df, 
                        file_path=file_path,
                        sheet_name=None if sheet_name=="__default__" else sheet_name
                    )
                )
        except Exception as e:
            self.logger.error(f"Error reading Table {file_path}: {e}")
        
        return results

    def _extract_pdf(self, file_path: Path) -> str:
        """Извлечение текста из pdf"""
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                txt = page.extract_text()
                if txt:
                    text_parts.append(txt)
        return "\n".join(text_parts)

    def _extract_docx(self, file_path: Path) -> str:
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _extract_text(self, file_path: Path) -> str:
        """функция для извлечения текста из файлов"""
        suffix = file_path.suffix.lower()
        if suffix == ".txt":    
            return file_path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".pdf":
            return self._extract_pdf(file_path)
        if suffix == ".docx":
            return self._extract_docx(file_path)
        if suffix == ".pptx":
            return self._extract_pptx(file_path)
        raise ValueError(f"Unsupported text format: {suffix}")   

    def _parse_faq_text(
        self,
        text: str,
        file_path: Path
    ) -> List[RawFAQItem]:

        results = []
        lines = [l.strip() for l in text.splitlines() if l.strip()]

        current_q = None
        current_a = []

        def flush():
            nonlocal current_q, current_a
            if current_q and current_a:
                results.append(
                    RawFAQItem(
                        question=current_q,
                        answer=" ".join(current_a),
                        section_path=[],
                        source_file=file_path.name,
                        source_type=file_path.suffix.lstrip(".")
                    )
                )
            current_q = None
            current_a = []

        for line in lines:
            if QUESTION_WORD_RE.search(line) or line.endswith("?"):
                flush()
                current_q = line
                continue

            if current_q:
                current_a.append(line)

        flush()
        return results

    # обработчик pptx 
    def _extract_pptx(self, file_path: Path) -> str:
        prs = Presentation(str(file_path))

        text_parts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        text_parts.append(text)

        return "\n".join(text_parts)

    # --- Orchestration ---
    def process_directory(self, input_dir: Path):
        input_dir = Path(input_dir)
        if not input_dir.exists():
            self.logger.error(f"Directory not found: {input_dir}")
            return

        raw_items = []
    
        # 1. Сбор данных
        for file in input_dir.rglob("*"):
            if file.name.startswith("~"): continue # Игнор временных файлов Excel
            suffix = file.suffix.lower()
            if suffix == '.md':
                self.logger.info(f"Processing MD: {file.name}")
                raw_items.extend(self._parse_markdown(file))
            elif suffix in ['.xlsx', '.xls', '.csv']:
                self.logger.info(f"Processing Table: {file.name}")
                raw_items.extend(self._parse_table(file))
            elif suffix in ['.txt', '.pdf', '.docx', '.pptx', '.doc']:
                self.logger.info(f"Processing Text FAQ: {file.name}")
                text = self._extract_text(file)
                raw_items.extend(self._parse_faq_text(text, file))

        # 2. Нормализация и Дедупликация
        for raw in raw_items:
            normalized = self._normalize_item(raw)
            if normalized:
                self.items.append(normalized)
            else:
                # Можно логировать пропущенные (дубли или пустые)
                pass

        self.logger.info(f"Processed {len(raw_items)} raw items into {len(self.items)} unique valid items.")

    @staticmethod
    def generate_content_hash(text: str, section_path: list[str], kb_id: str) -> str:
        canonical = FAQPreprocessor.canonical_text(text)
        context = "|".join([
            canonical,
            "/".join(section_path) if section_path else "",
            kb_id
        ])
        return hashlib.sha256(context.encode("utf-8")).hexdigest()

    def save(self):
        if not self.items:
            self.logger.warning("No items to save.")
            return
        
        self.output_file.parent.mkdir(parents=True, exist_ok=True)
        data = {"faqs": [asdict(item) for item in self.items]}
        
        try:
            with open(self.output_file, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            self.logger.info(f"Successfully saved to {self.output_file}")
        except Exception as e:
            self.logger.error(f"Failed to save JSON: {e}")

if __name__ == "__main__":
    INPUT_DIR = Path(r"")
    OUTPUT_FILE = Path("prepared/faq.normalized.json")

    processor = FAQPreprocessor(OUTPUT_FILE, log_dir=Path(r"nst-consultant-mcp-servers-pack\mcp-server-faq\faq_service"))
    processor.process_directory(INPUT_DIR)
    processor.save()
    print("Operation_complete")